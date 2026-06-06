"""
Генератор презентації PowerPoint для захисту диплому.
Тема: Розроблення навчальної платформи для набуття практичних навичок
       full-stack розробки на основі Monorepo-архітектури
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Configuration ---
OUTPUT_FILE = "presentation/diploma_presentation.pptx"
DIAGRAMS_DIR = "diagrams/output"

# Colors
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)    # Blue
COLOR_DARK = RGBColor(0x1F, 0x2A, 0x37)       # Dark
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)   # Light gray bg
COLOR_ACCENT = RGBColor(0x05, 0x96, 0x69)      # Green accent
COLOR_RED = RGBColor(0xDC, 0x26, 0x26)         # Red for problems


def create_presentation():
    """Create presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # 16:9
    prs.slide_height = Cm(19.05)
    return prs


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_DARK

    # Title
    left = Cm(3)
    top = Cm(5)
    width = Cm(28)
    height = Cm(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    if subtitle:
        top2 = Cm(10)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Cm(3))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

    return slide


def add_content_slide(prs, title, bullets, image_path=None):
    """Add a slide with title and bullet points, optionally with image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title bar
    left = Cm(0)
    top = Cm(0)
    width = Cm(33.867)
    height = Cm(3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Cm(2), Cm(0.5), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # Content area
    if image_path and os.path.exists(image_path):
        # Left side: bullets, Right side: image
        content_left = Cm(2)
        content_width = Cm(16)
        img_left = Cm(19)
        img_width = Cm(13)

        # Add image
        slide.shapes.add_picture(image_path, img_left, Cm(4), width=img_width)
    else:
        content_left = Cm(2)
        content_width = Cm(30)

    # Bullets
    content_top = Cm(4)
    content_height = Cm(14)
    txBox2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        p.space_after = Pt(8)

        if bullet.startswith("**") and bullet.endswith("**"):
            # Bold item
            run = p.add_run()
            run.text = bullet.strip("*")
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = COLOR_DARK
        else:
            run = p.add_run()
            run.text = "• " + bullet
            run.font.size = Pt(15)
            run.font.color.rgb = COLOR_DARK

    return slide


def add_comparison_slide(prs, title, data):
    """Add a slide with a comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Cm(2), Cm(0.5), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # Table
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Cm(2), Cm(4), Cm(30), Cm(13))
    table = table_shape.table

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if row_idx == 0:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_DARK
                p.font.color.rgb = COLOR_WHITE
            elif row_idx == len(data) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
                p.font.bold = True

    return slide


def add_image_slide(prs, title, image_path, subtitle=""):
    """Add a slide with a large centered image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Cm(2), Cm(0.5), Cm(30), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # Image centered
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Cm(5), Cm(4), width=Cm(24))

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Cm(2), Cm(17), Cm(30), Cm(2))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(12)
        run2.font.italic = True
        run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


# ============================================================
# SLIDES CONTENT
# ============================================================

def main():
    prs = create_presentation()

    # --- Slide 1: Title ---
    add_title_slide(
        prs,
        "Розроблення навчальної платформи\nдля набуття практичних навичок\nfull-stack розробки\nна основі Monorepo-архітектури",
        "Ключак Ольга Андріївна • КН-41д"
    )

    # --- Slide 2: Problem ---
    add_content_slide(prs, "Проблема", [
        "**Як набути реальний досвід без реального проєкту?**",
        "Frontend-студент — де взяти працюючий backend з API?",
        "Backend-студент — як перевірити API без готового клієнта?",
        "Існуючі курси: todo-app з mock-даними, без реальної взаємодії",
        "Немає структурованого шляху (що першим, що другим?)",
        "Ментор коштує $50-100/год і доступний обмежений час",
    ])

    # --- Slide 3: Existing solutions ---
    add_comparison_slide(prs, "Існуючі рішення не вирішують проблему", [
        ["Платформа", "Реальний проєкт", "Full-stack", "Спринти", "AI-ментор"],
        ["freeCodeCamp", "Ні", "Частково", "Ні", "Ні"],
        ["The Odin Project", "Так", "Частково", "Ні", "Ні"],
        ["Codecademy", "Ні", "Частково", "Ні", "Так"],
        ["Udemy/Coursera", "Частково", "Так", "Ні", "Ні"],
        ["Frontend Mentor", "Так", "Ні", "Ні", "Ні"],
        ["LeetCode/HackerRank", "Ні", "Ні", "Ні", "Ні"],
        ["Наша платформа", "Так", "Так", "Так", "Так"],
    ])

    # --- Slide 4: Our solution ---
    add_content_slide(prs, "Наше рішення", [
        "**Навчальна платформа на основі Monorepo**",
        "Студент обирає напрямок (frontend / backend / mobile)",
        "Запускає готові компоненти, реалізує свій проєкт по тікетах",
        "НЕ потрібно розуміти monorepo чи запускати вручну",
        "Розумна IDE (Kiro, Cursor) зчитує код → AI стає ментором",
        "Продукт = код + методика. AI — опціональний бонус",
        "Monorepo — технічне рішення під капотом",
    ], os.path.join(DIAGRAMS_DIR, "monorepo-structure.png"))

    # --- Slide 5: User scenario ---
    add_content_slide(prs, "Сценарій роботи студента", [
        "1. Відкриває IDE з 2 папками: платформа + свій проєкт",
        "2. AI бачить обидва контексти одночасно",
        '3. "Запусти backend" → AI запускає',
        '4. "Дай першу задачу" → AI видає тікет Sprint 1',
        "5. Студент реалізує задачу у своєму проєкті",
        '6. "Перевір мій код" → AI дає фідбек',
        "7. Переходить до наступної задачі",
        "",
        "Використання AI — опціонально",
    ], os.path.join(DIAGRAMS_DIR, "use-case.png"))

    # --- Slide 6: Sprint model ---
    add_content_slide(prs, "Спринтова модель навчання", [
        "**Sprint 1 — Авторизація**",
        "JWT, форми, валідація, захищені маршрути",
        "**Sprint 2 — Читання даних**",
        "API-запити, стани завантаження, обробка помилок",
        "**Sprint 3 — Повний CRUD**",
        "Створення, редагування, видалення, файли, пагінація",
        "**Sprint 4 — Розширення**",
        "Пріоритети, редизайн, адмін-панель",
        "",
        "Це реальні продакшн-сценарії = 90% щоденної роботи",
    ], os.path.join(DIAGRAMS_DIR, "sprint-progression.png"))

    # --- Slide 7: AI in IDE ---
    add_content_slide(prs, "Розумні IDE з вбудованим AI", [
        "Kiro, Cursor, Antigravity, GitHub Copilot — вже мають AI",
        "Відкриваєш платформу → AI зчитує контекст автоматично",
        "AI стає ментором БЕЗ додаткової розробки",
        "Студент спілкується як з людиною",
        "Опціонально — платформа працює і без AI",
        "Бонус: навички prompt engineering",
        "",
        '**"Розумна IDE перетворює код на ментора"**',
    ], os.path.join(DIAGRAMS_DIR, "ai-assistant-flow.png"))

    # --- Slide 8: Skills ---
    add_content_slide(prs, "Що отримує студент", [
        "Full-stack розробка (React / Express / React Native)",
        "Робота з реальним API та клієнт-серверною взаємодією",
        "TypeScript, Redux, REST API, JWT",
        "Робота з тікетами як у реальній компанії",
        "Git, npm, monorepo, сучасний інструментарій",
        "Система обробки помилок, валідація, безпека",
        "Prompt engineering та робота з AI",
        "Готовий проєкт у портфоліо",
    ])

    # --- Slide 9: Architecture ---
    add_image_slide(
        prs,
        "Архітектура серверного додатку",
        os.path.join(DIAGRAMS_DIR, "backend-layers.png"),
        "Express → Sprint Router → Auth Middleware → Controllers → Services → Sequelize → SQLite"
    )

    # --- Slide 10: Relevance ---
    add_content_slide(prs, "Актуальність у цифрах", [
        "38% зростання попиту на full-stack (LinkedIn 2025)",
        "82% розробників використовують AI щоденно (GitHub 2025)",
        "Зарплата full-stack в Україні: $3500-5500/міс (DOU 2026)",
        "Ментор: $50-100/год. IDE + платформа: безкоштовно, 24/7",
        "67% juniors: «немає практичного досвіду» = головна перешкода",
    ])

    # --- Slide 11: Demo ---
    add_content_slide(prs, "Демо: повний цикл студента", [
        "**1. Відкриття workspace** — IDE + 2 папки",
        '**2. Запуск** — "Запусти backend" → :4000 ready',
        '**3. Задача** — "Дай першу задачу Sprint 1"',
        '**4. Робота** — "Поясни що робити" → підказки',
        '**5. Перевірка** — "Оціни мій код" → фідбек',
        '**6. Далі** — "Яка наступна задача?"',
        "",
        "Весь процес — в одному вікні IDE",
    ])

    # --- Slide 12: Future ---
    add_content_slide(prs, "Подальший розвиток", [
        "Sprint 5-6: WebSocket, зовнішні API, RBAC",
        "Автоматизована перевірка коду (автотести для тікетів)",
        "Система прогресу та аналітики (дашборд)",
        "Підтримка інших стеків (NestJS, Vue, Flutter)",
        "Покращення AI-менторства (промпти per sprint)",
    ])

    # --- Slide 13: Conclusions ---
    add_content_slide(prs, "Висновки", [
        "✓ Створено платформу: 3 додатки + 5 packages + 85+ тікетів",
        "✓ Спринтова модель з реальними продакшн-сценаріями",
        "✓ Платформа дружня до розумних IDE (AI — автоматично)",
        "✓ Пілотне тестування: тікети 4.4/5, AI 4.8/5",
        "✓ Готова до використання у навчальних закладах",
        "",
        "**Дякую за увагу!**",
    ])

    # Save
    prs.save(OUTPUT_FILE)
    print(f"✅ Presentation saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
